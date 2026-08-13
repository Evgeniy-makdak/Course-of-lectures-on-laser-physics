# -*- coding: utf-8 -*-
"""Generate / extract figures and tables for the print handbook."""
from __future__ import annotations

import importlib.util
import io
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(r"C:\Users\Volkov\Desktop\Lecture-Unscheduled")
ASSETS = Path(__file__).resolve().parent / "assets"
ASSETS.mkdir(parents=True, exist_ok=True)


def _load(py_path: Path):
    name = py_path.stem + "_mod"
    spec = importlib.util.spec_from_file_location(name, py_path)
    mod = importlib.util.module_from_spec(spec)
    # Avoid running __main__
    sys.modules[name] = mod
    spec.loader.exec_module(mod)
    return mod


def _save(img, name: str) -> Path:
    path = ASSETS / name
    img.save(path, "PNG")
    print("fig", path.name)
    return path


def extract_pptx_pictures(pptx_path: Path, prefix: str) -> dict[int, Path]:
    """Return {1-based slide_index: png_path} for first picture on slide."""
    prs = Presentation(str(pptx_path))
    out = {}
    for i, slide in enumerate(prs.slides, 1):
        n = 0
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            n += 1
            blob = shape.image.blob
            ext = shape.image.ext or "png"
            path = ASSETS / f"{prefix}_s{i}_{n}.{ext}"
            path.write_bytes(blob)
            if i not in out:
                out[i] = path
            print("extract", path.name)
    return out


def prepare_all() -> dict:
    """Build asset catalog used by handbook builder."""
    catalog = {"figures": {}, "tables": {}}

    # Lecture 1 / Chapter 2 — graph from PPTX
    l1_pptx = next((ROOT / "Lecture-1-main").glob("Лекция_1_*.pptx"))
    pics1 = extract_pptx_pictures(l1_pptx, "ch2")
    if 11 in pics1:
        catalog["figures"]["ch2_threshold"] = {
            "path": pics1[11],
            "caption": "График пороговой плотности энергии для стали, титана и алюминия",
        }

    # Lecture 2 / Chapter 3
    m2 = _load(ROOT / "Lecture-2-main" / "build_lecture_2.py")
    figs2 = [
        ("ch3_spontaneous", m2.draw_spontaneous_transitions_pil,
         "Спонтанные переходы между уровнями E₂ и E₁"),
        ("ch3_avalanche", m2.draw_photon_avalanche_pil,
         "Фотонная лавина: вынужденное излучение и умножение фотонов"),
        ("ch3_inversion_prob", m2.draw_inversion_probability_pil,
         "Равенство вероятностей вынужденных переходов Ω₁₂ = Ω₂₁"),
        ("ch3_ruby", m2.draw_ruby_three_level_pil,
         "Трёхуровневая схема рубинового лазера (ионы Cr³⁺)"),
        ("ch3_pump", m2.draw_elliptical_pump_cavity_pil,
         "Эллиптический отражатель: лампа накачки и рубиновый стержень"),
        ("ch3_resonator", m2.draw_ruby_resonator_pil,
         "Оптический резонатор рубинового лазера"),
    ]
    for key, fn, cap in figs2:
        catalog["figures"][key] = {"path": _save(fn(), f"{key}.png"), "caption": cap}

    catalog["tables"]["ch3_laser_compare"] = {
        "caption": "Сравнение волоконного, диодного и CO₂-лазеров для аддитивных технологий",
        "headers": ["Параметр", "Волоконный", "Диодный", "CO₂"],
        "rows": [
            ["Длина волны λ", "≈ 1070 нм", "808–980 нм", "10,6 мкм"],
            ["КПД η", "30–40 %", "50–70 %", "10–20 %"],
            ["Качество пучка M²", "< 1,5", "> 10 (часто 20–30)", "1–1,5"],
            ["Поглощение Al", "≈ 40 % @ 1070 нм", "— (накачка)", "< 5 % (R > 95 %)"],
            ["Роль в СЛП", "Основной источник", "Накачка волоконных", "Редко; не для Al/Cu"],
        ],
    }

    # Lecture 3 / Chapter 1
    m3 = _load(ROOT / "Lecture-3-main" / "build_lecture_3.py")
    figs3 = [
        ("ch1_atom", m3.draw_atom_structure_pil,
         "Эволюция моделей атома и квантовые числа орбиталей"),
        ("ch1_pauli", m3.draw_pauli_shells_pil,
         "Принцип Паули и заполнение оболочек: N_max = 2n²"),
        ("ch1_sodium", m3.draw_sodium_example_pil,
         "Пример: переходы в атоме натрия на языке орбиталей и уровней"),
        ("ch1_two_schemes", m3.draw_two_diagrams_pil,
         "Две схемы: орбитали (пространство) и уровни энергии"),
        ("ch1_points", m3.draw_points_are_particles_pil,
         "Точки на линиях E₁, E₂ — частицы-излучатели, а не электроны"),
        ("ch1_not_always", m3.draw_not_always_orbital_pil,
         "Переход E₂→E₁ не всегда означает смену орбитали"),
    ]
    for key, fn, cap in figs3:
        catalog["figures"][key] = {"path": _save(fn(), f"{key}.png"), "caption": cap}

    catalog["tables"]["ch1_quantum"] = {
        "caption": "Квантовые числа электрона в атоме",
        "headers": ["Число", "Обозначение", "Смысл", "Возможные значения"],
        "rows": [
            ["Главное", "n", "Уровень / размер", "1, 2, 3, …"],
            ["Орбитальное", "l", "Форма орбитали", "0…n−1 (s, p, d, f, …)"],
            ["Магнитное", "m_l", "Ориентация", "−l … +l"],
            ["Спиновое", "m_s", "Спин электрона", "+1/2, −1/2"],
        ],
    }

    # Lecture 4 / Chapter 4
    m4 = _load(ROOT / "Lecture-4-main" / "build_lecture_4.py")
    figs4 = [
        ("ch4_energy", m4.draw_photon_energy_pil,
         "Энергия фотона и длина волны: 1070 нм и 355 нм"),
        ("ch4_spectrum", m4.draw_em_spectrum_pil,
         "Шкала электромагнитных излучений и технологические маркеры лазеров"),
        ("ch4_mechanisms", m4.draw_two_mechanisms_pil,
         "Два механизма: фотохимический (УФ) и тепловой (ИК)"),
        ("ch4_absorption", m4.draw_absorption_curves_pil,
         "Качественные кривые спектральной поглощательной способности A(λ)"),
        ("ch4_choice", m4.draw_laser_choice_pil,
         "Практические критерии выбора лазера и стратегии для алюминия"),
    ]
    for key, fn, cap in figs4:
        catalog["figures"][key] = {"path": _save(fn(), f"{key}.png"), "caption": cap}

    catalog["tables"]["ch4_absorption"] = {
        "caption": "Оценки поглощательной способности металлов на ключевых длинах волн",
        "headers": ["Материал", "A @ ≈1,07 мкм", "A @ 10,6 мкм", "Комментарий"],
        "rows": [
            ["Сталь", "≈ 30–50 %", "приемлемо в ряде задач", "Устойчивое СЛП на волоконном лазере"],
            ["Титан", "≈ 40–60 %", "—", "Высокое A + низкая теплопроводность"],
            ["Алюминий", "≈ 5–10 %", "≈ 2–3 %", "Низкое A + высокий теплоотвод"],
            ["Медь", "очень низкое", "ещё ниже", "Часто нужны green/blue лазеры"],
        ],
    }

    # Persist a simple manifest
    lines = ["# Assets", ""]
    for k, v in catalog["figures"].items():
        lines.append(f"- FIG {k}: {v['path'].name} — {v['caption']}")
    for k, v in catalog["tables"].items():
        lines.append(f"- TAB {k}: {v['caption']}")
    (ASSETS / "manifest.txt").write_text("\n".join(lines), encoding="utf-8")
    print("assets ready:", ASSETS)
    return catalog


if __name__ == "__main__":
    prepare_all()
